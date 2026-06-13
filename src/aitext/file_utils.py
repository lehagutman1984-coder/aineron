import uuid
import os
import base64
import tempfile
from pathlib import Path
from django.core.files.storage import default_storage
from django.core.files.base import ContentFile
import logging

logger = logging.getLogger(__name__)

# ========== ОСНОВНЫЕ ФУНКЦИИ ==========

def save_uploaded_file(file_data, message):
    """
    Сохраняет загруженный файл, связывая с сообщением.
    file_data: dict с полями name, type, dataUrl (из клиента)
    """
    try:
        from .models import FileAttachment

        original_name = file_data.get('name', 'file')
        data_url = file_data.get('dataUrl', '')
        mime = file_data.get('type', 'application/octet-stream')

        if not data_url:
            return None

        # Декодируем base64
        if data_url.startswith('data:'):
            header, data = data_url.split(',', 1)
            file_bytes = base64.b64decode(data)
        else:
            file_bytes = base64.b64decode(data_url)

        # Определяем расширение
        file_extension = Path(original_name).suffix.lower()
        if not file_extension:
            ext_map = {
                'image/jpeg': '.jpg',
                'image/png': '.png',
                'image/gif': '.gif',
                'image/webp': '.webp',
                'video/mp4': '.mp4',
                'application/pdf': '.pdf',
                'text/plain': '.txt',
                'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet': '.xlsx',
                'application/vnd.ms-excel': '.xls',
                'application/vnd.openxmlformats-officedocument.wordprocessingml.document': '.docx',
                'application/msword': '.doc',
                'application/vnd.openxmlformats-officedocument.presentationml.presentation': '.pptx',
                'application/vnd.ms-powerpoint': '.ppt',
            }
            file_extension = ext_map.get(mime, '.bin')

        # Генерируем уникальное имя для хранения
        storage_filename = f"attachments/{message.chat.user.id}/{message.chat.id}/{uuid.uuid4()}{file_extension}"
        file_obj = ContentFile(file_bytes, name=storage_filename)
        file_path = default_storage.save(storage_filename, file_obj)

        # Определяем тип медиа
        media_type = guess_media_type(mime, file_extension)

        # Создаём запись в БД
        attachment = FileAttachment(
            message=message,
            filename=original_name,
            file_path=file_path,
            file_size=len(file_bytes),
            mime_type=mime,
            media_type=media_type,
            source='uploaded'
        )

        # Сохраняем информацию (размеры изображения, если нужно)
        attachment.save_file_info(file_obj)
        attachment.save()

        # Если это текстовый/документный файл – извлекаем текст
        if media_type == 'other':
            try:
                with default_storage.open(file_path, 'rb') as f:
                    file_content = f.read()
                extracted = extract_text_from_file(None, original_name, file_content)
                if extracted:
                    attachment.extracted_text = extracted
                    attachment.save(update_fields=['extracted_text'])
                    logger.info(f"Извлечён текст из {original_name} ({len(extracted)} символов)")
            except Exception as e:
                logger.warning(f"Не удалось извлечь текст из {original_name}: {e}")

        return attachment

    except Exception as e:
        logger.error(f"Ошибка сохранения файла: {e}")
        return None


def guess_media_type(mime_type, extension):
    """Определяет тип медиа по MIME или расширению"""
    if mime_type.startswith('image/'):
        return 'image'
    if mime_type.startswith('video/'):
        return 'video'
    if mime_type.startswith('audio/'):
        return 'audio'
    if mime_type == 'application/pdf':
        return 'pdf'
    # По расширению
    image_exts = ['.jpg', '.jpeg', '.png', '.gif', '.webp', '.bmp', '.svg', '.ico', '.tiff']
    video_exts = ['.mp4', '.avi', '.mov', '.mkv', '.webm', '.flv', '.wmv', '.m4v', '.mpg', '.mpeg', '.3gp', '.ts']
    audio_exts = ['.mp3', '.wav', '.ogg', '.flac', '.m4a', '.aac', '.wma', '.opus', '.mid', '.midi']
    if extension in image_exts:
        return 'image'
    if extension in video_exts:
        return 'video'
    if extension in audio_exts:
        return 'audio'
    if extension == '.pdf':
        return 'pdf'
    return 'other'


def validate_file(file_data):
    """
    Валидация файла перед сохранением.
    Возвращает (is_valid, error_message, mime_type)
    """
    max_size = 50 * 1024 * 1024  # 50MB
    if file_data.get('size', 0) > max_size:
        return False, f"Файл слишком большой. Максимум 50MB", None

    name = file_data.get('name', '')
    mime = file_data.get('type', '')

    ext = Path(name).suffix.lower()

    # Если расширения нет, пробуем определить по MIME
    if not ext and mime:
        mime_to_ext = {
            'text/plain': '.txt',
            'text/html': '.html',
            'text/markdown': '.md',
            'application/pdf': '.pdf',
            'application/msword': '.doc',
            'application/vnd.openxmlformats-officedocument.wordprocessingml.document': '.docx',
            'application/vnd.ms-excel': '.xls',
            'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet': '.xlsx',
            'application/vnd.ms-powerpoint': '.ppt',
            'application/vnd.openxmlformats-officedocument.presentationml.presentation': '.pptx',
            'image/jpeg': '.jpg',
            'image/png': '.png',
            'image/gif': '.gif',
            'image/webp': '.webp',
            'video/mp4': '.mp4',
            'audio/mpeg': '.mp3',
            'application/zip': '.zip',
            'application/x-rar-compressed': '.rar',
            'application/x-7z-compressed': '.7z',
            'application/x-tar': '.tar',
            'application/gzip': '.gz',
        }
        ext = mime_to_ext.get(mime, '')
        if ext:
            name = name + ext

    allowed_extensions = {
        '.jpg', '.jpeg', '.png', '.gif', '.webp', '.bmp', '.svg', '.ico', '.tiff',
        '.mp4', '.avi', '.mov', '.mkv', '.webm', '.flv', '.wmv', '.m4v', '.mpg', '.mpeg', '.3gp', '.ts',
        '.mp3', '.wav', '.ogg', '.flac', '.m4a', '.aac', '.wma', '.opus', '.mid', '.midi',
        '.pdf',
        '.zip', '.rar', '.7z', '.tar', '.gz', '.bz2', '.tgz',
        '.txt', '.text', '.log', '.md', '.csv', '.tsv',
        '.html', '.htm', '.xml', '.json', '.yaml', '.yml',
        '.ini', '.cfg', '.conf', '.toml', '.properties',
        '.py', '.js', '.java', '.cpp', '.c', '.h', '.cs', '.php',
        '.rb', '.go', '.rs', '.swift', '.kt', '.sql', '.sh', '.bash',
        '.ps1', '.bat', '.cmd', '.lua', '.pl', '.r', '.ts', '.jsx', '.tsx',
        '.css', '.scss', '.sass', '.less', '.vue', '.ejs',
        '.rtf', '.eml', '.msg', '.ics', '.vcf',
        '.doc', '.docx', '.xls', '.xlsx', '.ppt', '.pptx', '.odt', '.ods', '.odp',
        '.epub', '.mobi', '.azw', '.chm'
    }

    if ext and ext not in allowed_extensions:
        return False, f"Неподдерживаемый формат файла: {ext}", None

    if not ext and mime:
        if (mime.startswith('text/') or mime.startswith('image/') or
            mime.startswith('video/') or mime.startswith('audio/') or
            mime in ['application/pdf', 'application/zip', 'application/x-rar-compressed',
                     'application/x-7z-compressed', 'application/x-tar', 'application/gzip',
                     'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet',
                     'application/vnd.ms-excel',
                     'application/vnd.openxmlformats-officedocument.wordprocessingml.document',
                     'application/msword',
                     'application/vnd.openxmlformats-officedocument.presentationml.presentation',
                     'application/vnd.ms-powerpoint']):
            return True, "OK", mime
        else:
            return False, f"Неподдерживаемый формат файла: неизвестный тип {mime}", None

    return True, "OK", mime


# ========== ИЗВЛЕЧЕНИЕ ТЕКСТА ==========

def extract_text_from_file(file_path, original_filename, file_data=None):
    """
    Универсальная функция извлечения текста из разных форматов.
    """
    file_extension = Path(original_filename).suffix.lower()
    temp_file_path = None
    try:
        if file_data is not None:
            with tempfile.NamedTemporaryFile(delete=False, suffix=file_extension) as tmp_file:
                tmp_file.write(file_data)
                temp_file_path = tmp_file.name
            file_to_process = temp_file_path
        else:
            file_to_process = file_path

        # Выбираем обработчик по расширению
        if file_extension in ['.zip', '.rar', '.7z', '.tar', '.gz', '.bz2', '.tgz']:
            return extract_text_from_archive(file_to_process, original_filename)
        elif file_extension in ['.xlsx', '.xls']:
            return extract_text_from_excel(file_to_process)
        elif file_extension in ['.docx', '.doc']:
            return extract_text_from_docx(file_to_process)
        elif file_extension in ['.pptx', '.ppt']:
            return extract_text_from_pptx(file_to_process)
        elif file_extension == '.pdf':
            return extract_text_from_pdf(file_to_process)
        else:
            return extract_text_from_txt(file_to_process)

    except Exception as e:
        logger.error(f"Ошибка извлечения текста из {original_filename}: {e}")
        return f"[НЕ УДАЛОСЬ ИЗВЛЕЧЬ ТЕКСТ ИЗ ФАЙЛА: {original_filename}]"
    finally:
        if temp_file_path and os.path.exists(temp_file_path):
            try:
                os.unlink(temp_file_path)
            except:
                pass


def extract_text_from_txt(file_path):
    """Чтение текстовых файлов с автоопределением кодировки"""
    encodings = ['utf-8', 'utf-8-sig', 'cp1251', 'windows-1251', 'koi8-r', 'latin-1']
    for enc in encodings:
        try:
            with open(file_path, 'r', encoding=enc) as f:
                text = f.read(20000)
                if text.strip():
                    logger.info(f"Текст извлечён (кодировка {enc}), длина {len(text)} символов")
                    return text
        except Exception:
            continue
    # Fallback
    with open(file_path, 'rb') as f:
        raw = f.read(20000)
        text = raw.decode('utf-8', errors='replace')
        logger.warning(f"Использована fallback кодировка, длина текста {len(text)}")
        return text


def extract_text_from_pdf(file_path):
    """Извлечение текста из PDF с помощью PyPDF2 или pdfplumber"""
    try:
        # Попробуем PyPDF2
        import PyPDF2
        text_parts = []
        with open(file_path, 'rb') as f:
            reader = PyPDF2.PdfReader(f)
            for page in reader.pages:
                page_text = page.extract_text()
                if page_text:
                    text_parts.append(page_text)
        full_text = "\n".join(text_parts)
        if full_text.strip():
            return full_text
        else:
            return "[PDF документ не содержит текста]"
    except ImportError:
        logger.warning("PyPDF2 не установлен, пробуем pdfplumber")
        try:
            import pdfplumber
            text_parts = []
            with pdfplumber.open(file_path) as pdf:
                for page in pdf.pages:
                    page_text = page.extract_text()
                    if page_text:
                        text_parts.append(page_text)
            full_text = "\n".join(text_parts)
            if full_text.strip():
                return full_text
            else:
                return "[PDF документ не содержит текста]"
        except ImportError:
            return "[PDF документ: установите PyPDF2 или pdfplumber]"
    except Exception as e:
        logger.error(f"Ошибка извлечения текста из PDF: {e}")
        return f"[Ошибка чтения PDF: {e}]"


def extract_text_from_excel(file_path):
    """Извлечение текста из Excel (.xlsx, .xls) с помощью openpyxl"""
    try:
        import openpyxl
        wb = openpyxl.load_workbook(file_path, read_only=True, data_only=True)
        text_parts = []
        for sheet_name in wb.sheetnames:
            sheet = wb[sheet_name]
            text_parts.append(f"=== Лист: {sheet_name} ===")
            # Собираем текст из всех ячеек, не пустых
            cell_texts = []
            for row in sheet.iter_rows(values_only=True):
                row_texts = [str(cell) for cell in row if cell is not None and str(cell).strip()]
                if row_texts:
                    cell_texts.append(" | ".join(row_texts))
            if cell_texts:
                text_parts.extend(cell_texts)
            else:
                text_parts.append("(пустой лист)")
        wb.close()
        full_text = "\n".join(text_parts)
        if full_text.strip():
            return full_text
        else:
            return "[Excel документ пуст]"
    except ImportError:
        return "[Excel документ: установите openpyxl]"
    except Exception as e:
        logger.error(f"Ошибка извлечения текста из Excel: {e}")
        return f"[Ошибка чтения Excel: {e}]"


def extract_text_from_docx(file_path):
    """Извлечение текста из DOCX с помощью python-docx"""
    try:
        from docx import Document
        doc = Document(file_path)
        text_parts = []
        for para in doc.paragraphs:
            if para.text.strip():
                text_parts.append(para.text)
        # Извлекаем текст из таблиц
        for table in doc.tables:
            for row in table.rows:
                row_text = []
                for cell in row.cells:
                    if cell.text.strip():
                        row_text.append(cell.text.strip())
                if row_text:
                    text_parts.append(" | ".join(row_text))
        full_text = "\n".join(text_parts)
        if full_text.strip():
            return full_text
        else:
            return "[Word документ пуст]"
    except ImportError:
        return "[Word документ: установите python-docx]"
    except Exception as e:
        logger.error(f"Ошибка извлечения текста из DOCX: {e}")
        return f"[Ошибка чтения DOCX: {e}]"


def extract_text_from_pptx(file_path):
    """Извлечение текста из PowerPoint с помощью python-pptx"""
    try:
        from pptx import Presentation
        prs = Presentation(file_path)
        text_parts = []
        for slide_num, slide in enumerate(prs.slides, 1):
            slide_text = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_text.append(shape.text.strip())
            if slide_text:
                text_parts.append(f"=== Слайд {slide_num} ===")
                text_parts.extend(slide_text)
        full_text = "\n".join(text_parts)
        if full_text.strip():
            return full_text
        else:
            return "[PowerPoint презентация пуста]"
    except ImportError:
        return "[PowerPoint презентация: установите python-pptx]"
    except Exception as e:
        logger.error(f"Ошибка извлечения текста из PPTX: {e}")
        return f"[Ошибка чтения PPTX: {e}]"


def extract_text_from_archive(file_path, original_filename):
    """Извлечение текста из архивов (распаковка и чтение содержимого)"""
    try:
        import zipfile
        import tarfile
        import tempfile

        text_parts = [f"=== АРХИВ: {original_filename} ==="]

        with tempfile.TemporaryDirectory() as temp_dir:
            if str(file_path).lower().endswith('.zip'):
                with zipfile.ZipFile(file_path, 'r') as zip_ref:
                    zip_ref.extractall(temp_dir)
                    file_names = [f for f in zip_ref.namelist() if not f.endswith('/')]
            elif any(str(file_path).lower().endswith(ext) for ext in ['.tar', '.tar.gz', '.tgz', '.tar.bz2']):
                mode = 'r'
                if str(file_path).lower().endswith(('.tar.gz', '.tgz')):
                    mode = 'r:gz'
                elif str(file_path).lower().endswith('.tar.bz2'):
                    mode = 'r:bz2'
                with tarfile.open(file_path, mode) as tar_ref:
                    tar_ref.extractall(temp_dir)
                    file_names = [m.name for m in tar_ref.getmembers() if m.isfile()]
            else:
                text_parts.append("Неподдерживаемый формат архива")
                return "\n".join(text_parts)

            if not file_names:
                text_parts.append("Архив пуст")
                return "\n".join(text_parts)

            # Ограничиваем количество файлов
            file_names = file_names[:50]
            text_parts.append(f"Количество файлов: {len(file_names)}")
            text_parts.append("\nСОДЕРЖИМОЕ:\n" + "=" * 60)

            for fname in file_names:
                full_path = os.path.join(temp_dir, fname)
                if os.path.isfile(full_path):
                    ext = Path(fname).suffix.lower()
                    # Пропускаем бинарные файлы (изображения, видео, аудио, исполняемые, архивы)
                    skip_exts = ['.exe', '.dll', '.so', '.bin', '.jpg', '.jpeg', '.png', '.gif', '.bmp', '.tiff',
                                 '.ico', '.svg', '.webp', '.mp3', '.mp4', '.avi', '.mov', '.mkv', '.wav', '.ogg',
                                 '.flac', '.pdf', '.zip', '.rar', '.7z', '.tar', '.gz', '.bz2']
                    if ext in skip_exts:
                        continue
                    try:
                        content = extract_text_from_txt(full_path)
                        if content and content.strip():
                            text_parts.append(f"\n--- {fname} ---")
                            if len(content) > 3000:
                                content = content[:6000] + "\n[... обрезано ...]"
                            text_parts.append(content)
                    except:
                        pass
        return "\n".join(text_parts)
    except Exception as e:
        logger.error(f"Ошибка обработки архива {original_filename}: {e}")
        return f"[АРХИВ: {original_filename} - ошибка обработки]"


# ========== ПОДГОТОВКА ДЛЯ AI ==========

def prepare_media_for_ai(file_attachment):
    """
    Подготавливает файл для отправки в OpenRouter.
    Для изображений – возвращает data URL.
    Для остальных – извлечённый текст (если есть) или описание.
    """
    try:
        if file_attachment.media_type == 'image':
            with default_storage.open(file_attachment.file_path, 'rb') as f:
                file_bytes = f.read()
            b64 = base64.b64encode(file_bytes).decode('utf-8')
            mime = file_attachment.mime_type or 'image/jpeg'
            return {
                'type': 'image_url',
                'image_url': {
                    'url': f"data:{mime};base64,{b64}"
                }
            }
        else:
            # Для текстовых файлов и документов
            if file_attachment.extracted_text:
                text = file_attachment.extracted_text
                logger.info(f"Для файла {file_attachment.filename} передан текст длиной {len(text)} символов")
            else:
                text = file_attachment.get_file_description()
                logger.warning(f"Для файла {file_attachment.filename} нет извлечённого текста, отправляю описание")
            return {
                'type': 'text',
                'text': text[:6000]  # ограничим длину
            }
    except Exception as e:
        logger.error(f"Ошибка подготовки файла {file_attachment.id} для AI: {e}")
        return {
            'type': 'text',
            'text': f"[Ошибка обработки файла: {file_attachment.filename}]"
        }


# ========== ДОПОЛНИТЕЛЬНЫЕ ФУНКЦИИ ==========

def get_files_for_message(message):
    """Получает все файлы, прикреплённые к сообщению"""
    try:
        from .models import FileAttachment
        attachments = message.attachments.all()
        file_data = []
        for att in attachments:
            file_data.append({
                'id': str(att.id),
                'filename': att.filename,
                'url': att.file_url,
                'size': att.file_size,
                'mime_type': att.mime_type,
                'media_type': att.media_type,
                'dimensions': att.dimensions,
                'duration': att.duration,
                'is_generated': att.is_generated,
                'created_at': att.created_at.isoformat() if att.created_at else None,
                'extracted_text': att.extracted_text if att.media_type == 'other' else None
            })
        return file_data
    except Exception as e:
        logger.error(f"Ошибка получения файлов сообщения: {e}")
        return []


def delete_file_attachment(file_attachment):
    """Удаляет файл из хранилища и БД"""
    try:
        if default_storage.exists(file_attachment.file_path):
            default_storage.delete(file_attachment.file_path)
        file_attachment.delete()
        return True
    except Exception as e:
        logger.error(f"Ошибка удаления файла {file_attachment.id}: {e}")
        return False


def get_file_statistics(user=None):
    """Возвращает статистику по файлам пользователя (или общую)"""
    try:
        from .models import FileAttachment
        from django.db.models import Sum

        queryset = FileAttachment.objects.all()
        if user:
            queryset = queryset.filter(message__chat__user=user)

        total_files = queryset.count()
        uploaded_files = queryset.filter(source='uploaded').count()
        generated_files = queryset.filter(source='ai_generated').count()

        images_count = queryset.filter(media_type='image').count()
        videos_count = queryset.filter(media_type='video').count()
        audio_count = queryset.filter(media_type='audio').count()
        pdf_count = queryset.filter(media_type='pdf').count()
        other_count = queryset.filter(media_type='other').count()

        total_size = queryset.aggregate(Sum('file_size'))['file_size__sum'] or 0

        type_stats = {}
        for file in queryset:
            mime = file.mime_type or 'unknown'
            type_stats[mime] = type_stats.get(mime, 0) + 1

        return {
            'total_files': total_files,
            'uploaded_files': uploaded_files,
            'generated_files': generated_files,
            'images_count': images_count,
            'videos_count': videos_count,
            'audio_count': audio_count,
            'pdf_count': pdf_count,
            'other_count': other_count,
            'total_size_bytes': total_size,
            'total_size_mb': round(total_size / (1024 * 1024), 2),
            'type_stats': type_stats
        }
    except Exception as e:
        logger.error(f"Ошибка получения статистики файлов: {e}")
        return {}
